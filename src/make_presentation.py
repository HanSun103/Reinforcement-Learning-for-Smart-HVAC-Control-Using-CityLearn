"""
make_presentation.py
--------------------
Generate presentation-ready PowerPoint decks from repo outputs.

Usage:
    python src/make_presentation.py
"""

from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
OUT_PPTX = ROOT / "HVAC_RL_Class_Presentation.pptx"
OUT_PPTX_V2 = ROOT / "HVAC_RL_Class_Presentation_v2.pptx"


def load_kpis(path: Path) -> Dict[str, float]:
    """Load Building_1 KPIs from a CSV into a metric->value dict."""
    df = pd.read_csv(path)
    out = {}
    for _, row in df.iterrows():
        metric = str(row["cost_function"])
        value = row.get("Building_1")
        if pd.notna(value):
            out[metric] = float(value)
    return out


def fmt(v: float) -> str:
    return f"{v:.4f}"


def add_notes(slide, notes: Optional[str]) -> None:
    """Attach presenter notes to a slide."""
    if not notes:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def add_title_and_bullets(
    prs: Presentation,
    title: str,
    bullets: List[str],
    notes: Optional[str] = None,
    bullet_font_size: int = 20,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(bullet_font_size)

    add_notes(slide, notes)


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str,
    notes: Optional[str] = None,
    image_top: float = 1.3,
    image_width: float = 12.0,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(image_top), width=Inches(image_width))
    else:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(1.5))
        box.text_frame.text = f"Image not found: {image_path.name}"

    cap = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(12.0), Inches(0.6))
    cap_tf = cap.text_frame
    cap_tf.text = caption
    cap_tf.paragraphs[0].font.size = Pt(16)
    add_notes(slide, notes)


def add_two_image_slide(
    prs: Presentation,
    title: str,
    left_img: Path,
    right_img: Path,
    left_cap: str,
    right_cap: str,
    notes: Optional[str] = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    if left_img.exists():
        slide.shapes.add_picture(str(left_img), Inches(0.4), Inches(1.2), width=Inches(6.1))
    if right_img.exists():
        slide.shapes.add_picture(str(right_img), Inches(6.8), Inches(1.2), width=Inches(6.1))

    cap_l = slide.shapes.add_textbox(Inches(0.5), Inches(6.45), Inches(6.0), Inches(0.5))
    cap_l.text_frame.text = left_cap
    cap_l.text_frame.paragraphs[0].font.size = Pt(14)

    cap_r = slide.shapes.add_textbox(Inches(6.9), Inches(6.45), Inches(6.0), Inches(0.5))
    cap_r.text_frame.text = right_cap
    cap_r.text_frame.paragraphs[0].font.size = Pt(14)
    add_notes(slide, notes)


def make_base_deck(prs: Presentation, compact: bool) -> None:
    """Create the shared deck content with optional compact style."""
    baseline = load_kpis(RESULTS / "baseline_kpis.csv")
    sac = load_kpis(RESULTS / "sac_kpis.csv")
    ppo = load_kpis(RESULTS / "ppo_kpis.csv")
    td3 = load_kpis(RESULTS / "td3_kpis.csv")

    add_title_and_bullets(
        prs,
        "Reinforcement Learning for Smart HVAC Control Using CityLearn",
        [
            "Course project prototype using CityLearn 2023 local evaluation dataset.",
            "Controllers compared: Rule-Based Controller (RBC), SAC, PPO, TD3.",
            "Objective: reduce electricity, cost, and carbon while preserving comfort.",
            "This deck reports only implemented pipeline and saved repo results.",
        ],
        notes="Open with scope: this is a real implemented prototype, not a future concept.",
        bullet_font_size=20 if compact else 22,
    )

    add_title_and_bullets(
        prs,
        "Problem Formulation (MDP)",
        [
            "Environment: CityLearn central agent; default single-building (Building_1).",
            "State: runtime-activated 23 observations including weather, load, SOC, pricing, carbon.",
            "Action: 3-D continuous HVAC/storage control vector through SB3 wrappers.",
            "Reward: -(0.4*energy + 0.4*comfort + 0.2*carbon) from src/reward.py.",
            "Episode length in runs: ~719 hourly steps.",
        ],
        notes="Mention 23 is from runtime activation logs, not just static config comments.",
        bullet_font_size=20 if compact else 19,
    )

    add_title_and_bullets(
        prs,
        "Methods and Evaluation Setup",
        [
            "RBC baseline: time-of-use charge/discharge schedule.",
            "SAC: off-policy stochastic; PPO: on-policy stochastic; TD3: off-policy deterministic.",
            "Train script runs baseline + selected RL algorithms; evaluate script compares all available models.",
            "EvalCallback saves best checkpoint per algorithm (results/best_{algo}/best_model.zip).",
            "KPIs and plots are generated automatically into results/.",
        ],
        notes="Frame method choice as representative, not exhaustive benchmarking.",
        bullet_font_size=20 if compact else 18,
    )

    add_image_slide(
        prs,
        "KPI Comparison Across Controllers",
        RESULTS / "kpi_comparison.png",
        "Lower is better for most normalized KPIs (reference = 1.0).",
        notes="Focus on electricity, cost, and carbon first.",
        image_top=1.3 if compact else 1.15,
    )

    add_title_and_bullets(
        prs,
        "Key KPI Numbers (from saved CSVs)",
        [
            f"Electricity total: RBC {fmt(baseline['electricity_consumption_total'])}, "
            f"SAC {fmt(sac['electricity_consumption_total'])}, "
            f"PPO {fmt(ppo['electricity_consumption_total'])}, "
            f"TD3 {fmt(td3['electricity_consumption_total'])}.",
            f"Cost total: RBC {fmt(baseline['cost_total'])}, SAC {fmt(sac['cost_total'])}, "
            f"PPO {fmt(ppo['cost_total'])}, TD3 {fmt(td3['cost_total'])}.",
            f"Carbon total: RBC {fmt(baseline['carbon_emissions_total'])}, "
            f"SAC {fmt(sac['carbon_emissions_total'])}, "
            f"PPO {fmt(ppo['carbon_emissions_total'])}, "
            f"TD3 {fmt(td3['carbon_emissions_total'])}.",
            "Comfort trade-off: SAC discomfort 0.5792 vs PPO/TD3 0.9804.",
        ],
        notes="State one clear takeaway: better energy/carbon can conflict with comfort.",
        bullet_font_size=20 if compact else 18,
    )

    add_two_image_slide(
        prs,
        "Behavior and Training Dynamics",
        RESULTS / "training_rewards.png",
        RESULTS / "reward_comparison.png",
        "Training rewards from monitor logs",
        "Per-step reward comparison (smoothed)",
        notes="Emphasize this is a short-horizon prototype, not final convergence.",
    )

    add_two_image_slide(
        prs,
        "Comfort Trade-off Evidence",
        RESULTS / "temperature_trace_sac.png",
        RESULTS / "temperature_trace_ppo.png",
        "SAC temperature trace",
        "PPO temperature trace",
        notes="Use this visual to explain why comfort KPIs differ.",
    )

    add_title_and_bullets(
        prs,
        "Business Plan (Realistic Productization View)",
        [
            "Likely users: facility operators, campus energy teams, BMS vendors.",
            "Value proposition: adaptive control can lower cost/carbon beyond fixed schedules.",
            "Needed before deployment: BAS integration, forecast pipeline, data QA, safety constraints.",
            "RBC still matters as fallback control and trust-building benchmark in operations.",
            "Deployment risks: comfort violations, model drift, robustness and cyber/security concerns.",
        ],
        notes="Keep this grounded: practical integration and safety before rollout.",
        bullet_font_size=20 if compact else 18,
    )

    add_title_and_bullets(
        prs,
        "Assumptions and Limitations",
        [
            "Validation is simulation-only (CityLearn), not live building deployment.",
            "Reported results are single-building default mode (MULTI_BUILDING=False).",
            "Short training budget (5 episodes) prioritizes prototype speed over full convergence.",
            "Some configured observations are unavailable in schema and skipped at runtime.",
            "Optuna tuning exists, but tuned params must be copied manually into config for full runs.",
        ],
        notes="Be transparent here; this usually improves Q&A outcomes.",
        bullet_font_size=20 if compact else 18,
    )

    add_title_and_bullets(
        prs,
        "Conclusion and Next Steps",
        [
            "Implemented end-to-end pipeline: tune/train/evaluate with reproducible outputs.",
            "RL methods clearly improved energy-cost-carbon KPIs vs RBC in current experiments.",
            "Comfort remains the key unresolved trade-off and should be constraint-prioritized next.",
            "Next phase: comfort-aware reward/constraints, multi-seed tests, multi-building experiments.",
            "Future deployment path requires integration, safety guardrails, and staged rollout.",
        ],
        notes="Close with balanced claim: promising results plus clear next steps.",
        bullet_font_size=20 if compact else 18,
    )


def build_presentation() -> None:
    """Original dense deck."""
    prs = Presentation()
    make_base_deck(prs, compact=True)
    prs.save(OUT_PPTX)
    print("Presentation created successfully.")


def build_presentation_v2() -> None:
    """Cleaner deck with larger visuals and embedded speaker notes."""
    prs = Presentation()
    make_base_deck(prs, compact=False)
    prs.save(OUT_PPTX_V2)
    print("Presentation v2 created successfully.")


if __name__ == "__main__":
    build_presentation()
    build_presentation_v2()
